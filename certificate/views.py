from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from .models import *
from .convert import ppt2pdf
import collections 
import collections.abc
from pptx import Presentation
from django.core.mail import send_mail, EmailMessage
import requests
import os
import pandas as pd
import time
import hashlib
# Create your views here.


def verify(request, verify_id):
	if ParticipantCertificate.objects.filter(certificatehash=verify_id).exists():
		certid =  ParticipantCertificate.objects.get(certificatehash=verify_id).certificateid
		return render(request,'cert_verify.html', {'certid' : certid})
	else:
		return render(request,'404.html')



def create(request):
	if request.user.is_staff and CertificateManager.objects.filter(email=request.user.email).exists():
		if request.method == "POST":
			csv = request.FILES.get('csv')
			temp = request.FILES.get('template')
			event_id = request.POST.get('event_id')
			event_name = Event.objects.get(event_id=event_id).event_name
			event = EventCertificate(
				user = request.user.email,
				event_id = event_id,
				event_name = event_name,
				data_file = csv,
				template = temp)
			event.save()
			return redirect(f"/certificate/{event.id}/{event.slug}")
		else:
			events = Event.objects.all().order_by('-event_date').values()
			return render(request, 'cert_create_event.html' , { 'events' : events})
	else:
		messages.success(request, 'Your are not Authorised')
		return redirect('home')
	
	
	
def track(request, id, slug):
	if request.user.is_staff and CertificateManager.objects.filter(email=request.user.email).exists():
		event = EventCertificate.objects.filter(slug=slug, id=id).first()
  
		if event.message:
			return render(request, 'cert_track.html', {
				'event_name': event.event_name,
				'event_date': event.date,
				'participat_details': ParticipantCertificate.objects.filter(event=event.id)
				})
   
		prs = Presentation(event.template)
		st=""
		for slide in prs.slides:
			for shape in slide.shapes:
				if shape.has_text_frame:
					st = st + shape.text
					st = st + " "
		li = st.split()
		tags = []
		for i in li:
			if i[0] == "<" and i[-1] == ">":
				tags.append(i)
		
		if request.method == "POST":
			email_col = request.POST.get('emails')
			subject = request.POST.get('subject')
			mess = request.POST.get('mess')
			values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]	
			print(values)
			event.email_column = email_col
			event.message = mess
			event.subject = subject
			event.save()
			
			df=pd.read_csv(event.data_file)
			df_len=df.shape
			i=0

			while i < df_len[0]:
				prs = Presentation(event.template)
				j=""
				if i<9:
					j="00"
				elif i>=9 and i < 99 :
					j="0"
				hash = ""
				url = ""
				for tag, v_type, value in values:
					for slide in prs.slides:
						for shape in slide.shapes:
							if shape.has_text_frame:
								if(shape.text.find(tag))!=-1:
										text_frame = shape.text_frame
										for paragraph in text_frame.paragraphs:
											for run in paragraph.runs:
												cur_text = run.text
												if v_type == 'text':
													new_text = cur_text.replace(tag, value)
												elif v_type == 'date':
													new_text = cur_text.replace(tag, '/'.join(value.split('-')[::-1]))
												elif v_type == 'csv':
													new_text = cur_text.replace(tag, df.loc[i,value])
												elif v_type == 'verifylink':
													str2hash = str(time.time())+ str(j)+str(i+1)
													hash = hashlib.md5(str2hash.encode())
													url = 'https://cysec.gitam.edu/certificate/verify/'+ str(hash.hexdigest())
													new_text = cur_text.replace(tag, url )
													print(new_text)
												elif v_type == "auto":
													new_text = cur_text.replace(tag, value+"/"+j+str(i+1))
												else:
													pass
												run.text = new_text
												
				
				s_name = df.loc[i,event.email_column].split('@')[0]

				prs.save(s_name+".pptx")
				parent_id = Event.objects.get(event_id=event.event_id).parentfolderid
				f_id = ppt2pdf(s_name+".pptx",s_name, parent_id)
				r = requests.get(f"https://docs.google.com/presentation/d/{f_id}/export/pdf", allow_redirects=True)
				open(s_name+'.pdf', 'wb').write(r.content)

				try:
					mail = EmailMessage(subject,
						f"Hello, {s_name} \n{mess}",
						settings.EMAIL_HOST_USER,
						[df.loc[i,event.email_column]])
					mail.attach_file(s_name+'.pdf')
					mail.send()
					ParticipantCertificate(
		 					event=event.id, 
			  				eventid=event.event_id,  
				  			email=df.loc[i,event.email_column], 
							certificatelink= url,
							certificatehash = hash.hexdigest(),
							certificateid = f_id,
					 		status=True
							).save()
					os.remove(s_name+'.pdf')
					os.remove(s_name+".pptx")
				except:
					ParticipantCertificate(
		 					event=event.id, 
			  				eventid=event.event_id,  
				  			email=df.loc[i,event.email_column], 
							certificatelink= url,
							certificatehash = hash.hexdigest(),
							certificateid = f_id,
					 		status=False
							).save()
					os.remove(s_name+'.pdf')
					os.remove(s_name+".pptx")
				i=i+1

			messages.success(request, "Certificates Sent Successfuly !!")
			return redirect(f"/certificate/{event.id}/{event.slug}")


		return render(request, 'cert_map_tags.html',{
			'columns': list(pd.read_csv(event.data_file).columns),
			'tags': tags,
			})
	else:
		messages.success(request, 'Your are not Authorised')
		return redirect('home')
  
  
def view_certificate_status(request):
	if request.user.is_staff and CertificateManager.objects.filter(email=request.user.email).exists():
		return render(request, 'cert_view_status.html',{
		'events': EventCertificate.objects.filter(user=request.user.email)
		})
	else:
		messages.success(request, 'Your are not Authorised')
		return redirect('home')
 
 
 
# def delete_event(request, id, slug):
# 	if request.user.is_staff and CertificateManager.objects.filter(email=request.user.email).exists():
# 		event = EventCertificate.objects.filter(slug=slug, id=id).first()
# 		if event.user == request.user.email:
# 			event.delete()
# 		return redirect('view_certificate_status')
# 	else:
# 		messages.success(request, 'Your are not Authorised')
# 		return redirect('home')